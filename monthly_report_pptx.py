"""
台中市房地產月報生成器（設計師版模板）
執行方式：python monthly_report_pptx.py

策略：開啟設計師格式化模板，更新文字與圖表資料，保留所有視覺設計。
Step 1：填入數據（圖表 + KPI）
Step 2：Gemini AI 自動產出文字摘要
"""

import sys, subprocess, re, json

for pkg, imp in [
    ("requests","requests"), ("pandas","pandas"),
    ("python-pptx","pptx"), ("gspread","gspread"),
    ("google-genai","google.genai"),
]:
    try:
        __import__(imp)
    except ImportError:
        print(f"安裝 {pkg}...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", pkg, "-q"])

import io, requests, pandas as pd
import gspread
from google import genai
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData, BubbleChartData

# ── 設定 ──────────────────────────────────────────────────────
SHEET_ID             = "1pN9_h5Pqe6CewXs8WPULSNpW8tXUKj1h8nZgMneu4HE"
GID_DETAIL           = 501026080    # 預售屋總表（14欄月份格式）
OUTPUT_DIR           = Path("/Volumes/ADMM/1-2公司基礎設備設定/實價登錄報告/報表產出")
TEMPLATE             = Path("/Volumes/ADMM/1-2公司基礎設備設定/實價登錄報告/報表模板/template_月報.pptx")
SERVICE_ACCOUNT_FILE = Path.home() / ".config/gspread/service_account.json"
GEMINI_KEY_FILE      = Path.home() / ".config/gemini_api_key.txt"
ZH_FONT              = "PingFang TC"

# 設計師配色
C_GOLD  = RGBColor(0xCF, 0xAD, 0x1E)
C_RED   = RGBColor(0xE7, 0x4C, 0x3C)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_GRAY  = RGBColor(0x95, 0xA5, 0xA6)

if not TEMPLATE.exists():
    print(f"找不到模板：{TEMPLATE}"); sys.exit(1)
if not SERVICE_ACCOUNT_FILE.exists():
    print(f"找不到 Service Account 金鑰：{SERVICE_ACCOUNT_FILE}"); sys.exit(1)
if not GEMINI_KEY_FILE.exists():
    print(f"找不到 Gemini API Key：{GEMINI_KEY_FILE}"); sys.exit(1)

# ── Gemini 初始化 ──────────────────────────────────────────────
_gc = genai.Client(api_key=GEMINI_KEY_FILE.read_text(encoding="utf-8").strip())
GEMINI_MODEL = "gemini-2.0-flash-lite"

def ask_gemini(prompt):
    try:
        resp = _gc.models.generate_content(model=GEMINI_MODEL, contents=prompt)
        return resp.text.strip()
    except Exception as e:
        print(f"  Gemini 呼叫失敗：{e}")
        return ""

# ── 讀資料（gspread Service Account）────────────────────────
gc = gspread.service_account(filename=str(SERVICE_ACCOUNT_FILE))

def fetch_sheet(gid):
    sh = gc.open_by_key(SHEET_ID)
    ws = next((w for w in sh.worksheets() if w.id == gid), None)
    if ws is None:
        raise ValueError(f"找不到 gid={gid} 的工作表")
    rows = ws.get_all_values()
    if len(rows) < 2:
        raise ValueError("工作表資料為空")
    return pd.DataFrame(rows[1:], columns=rows[0])

print("讀取總表...")
df = fetch_sheet(GID_DETAIL)
df.columns = ["id","成交年月","縣市","鄉鎮市區","交易標的","建案名稱","門牌",
              "建物型態","總價","單價","面積","屋齡","交易年月日","匯入時間"]
df["總價"] = pd.to_numeric(df["總價"], errors="coerce")
df["單價"] = pd.to_numeric(df["單價"], errors="coerce")

# ── 期別 ────────────────────────────────────────────────────
def m_key(m):
    return int(str(m).replace("-", ""))

months   = sorted(df["成交年月"].dropna().unique(), key=m_key, reverse=True)
latest_m = months[0]
prev_m   = months[1] if len(months) > 1 else None
now      = datetime.now()

df_l = df[df["成交年月"] == latest_m].copy()

# ── 行政區聚合 ──────────────────────────────────────────────
print("計算行政區聚合...")
ds = (df_l.groupby("鄉鎮市區", as_index=False)
      .agg(
          成交筆數=("id",   "count"),
          均單    =("單價", "mean"),
          中位單  =("單價", "median"),
          最高單  =("單價", "max"),
          最低單  =("單價", "min"),
          均總    =("總價", "mean"),
          最高總  =("總價", "max"),
          最低總  =("總價", "min"),
      )
      .sort_values("成交筆數", ascending=False)
      .reset_index(drop=True))
ds.rename(columns={"鄉鎮市區": "行政區"}, inplace=True)
for c in ["均單","中位單","最高單","最低單","均總","最高總","最低總"]:
    ds[c] = ds[c].round(2)

# ── KPI ──────────────────────────────────────────────────────
total_tx = len(df_l)
avg_unit = round(df_l["單價"].mean(), 2)
top_dist = ds.iloc[0]["行政區"] if not ds.empty else "-"
top_tx   = int(ds.iloc[0]["成交筆數"]) if not ds.empty else 0
num_proj = df_l["建案名稱"].nunique()
num_dist = len(ds)

if prev_m:
    df_p   = df[df["成交年月"] == prev_m]
    tx_chg = total_tx - len(df_p)
    tx_pct = round(tx_chg / len(df_p) * 100, 1) if len(df_p) else 0
    u_chg  = round(avg_unit - df_p["單價"].mean(), 2)
    u_pct  = round(u_chg / df_p["單價"].mean() * 100, 1) if df_p["單價"].mean() else 0
else:
    tx_chg = tx_pct = u_chg = u_pct = 0

def chg_str(val, pct, unit=""):
    arrow = "▲" if val > 0 else "▼" if val < 0 else "─"
    sign  = "+" if val > 0 else ""
    return f"較上月 {arrow} {sign}{val}{unit}（{sign}{pct}%）"

def chg_color(val):
    return C_RED if val > 0 else C_GREEN if val < 0 else C_GRAY

# ── 排行榜資料（供 Gemini 使用）──────────────────────────────
rank_cnt = (df_l.groupby("建案名稱")
            .agg(n=("id","count"), 區=("鄉鎮市區","first"))
            .reset_index().sort_values("n", ascending=False).head(10))
rank_unit = (df_l.groupby("建案名稱")
             .agg(v=("單價","max"), 區=("鄉鎮市區","first"))
             .reset_index().sort_values("v", ascending=False).head(10))
rank_tot = (df_l.groupby("建案名稱")
            .agg(v=("總價","max"), 區=("鄉鎮市區","first"))
            .reset_index().sort_values("v", ascending=False).head(10))
type_cnt = df_l["建物型態"].value_counts().head(6)

top5       = ds.head(5)["行政區"].tolist()
months_asc = sorted(df["成交年月"].dropna().unique(), key=m_key)

# ══════════════════════════════════════════════════════════════
# Gemini AI 文字生成
# ══════════════════════════════════════════════════════════════
print("AI 生成文字摘要...")

_base = f"""你是台中市預售屋市場分析師，用繁體中文撰寫，語氣專業簡潔，數字直接引用資料。
統計月份：{latest_m}"""

# ── Slide 1：三大觀察 ──────────────────────────────────────
_s1_data = f"""
總成交筆數：{total_tx}件（較上月 {tx_chg:+d}件，{tx_pct:+.1f}%）
全市均單：{avg_unit}萬/坪（較上月 {u_chg:+.2f}萬，{u_pct:+.1f}%）
成交量第一行政區：{top_dist}（{top_tx}件）
活躍建案數：{num_proj}個，涵蓋行政區：{num_dist}區
成交前5區：{', '.join(f"{r['行政區']}({int(r['成交筆數'])}件)" for _, r in ds.head(5).iterrows())}
"""
_s1_prompt = f"""{_base}
{_s1_data}
請輸出3條市場觀察，格式嚴格如下（共6行，不要其他內容）：
01)標題（5字以內）
說明（25字以內，含具體數字）
02)標題（5字以內）
說明（25字以內，含具體數字）
03)標題（5字以內）
說明（25字以內，含具體數字）"""

_s1_raw = ask_gemini(_s1_prompt)
_s1_lines = [l.strip() for l in _s1_raw.splitlines() if l.strip()]
def _extract_obs(lines, idx):
    """從輸出取第 idx 條（0-based），回傳 '標題\n說明' 格式。"""
    title = lines[idx*2]   if idx*2   < len(lines) else f"0{idx+1})觀察{idx+1}"
    desc  = lines[idx*2+1] if idx*2+1 < len(lines) else "（待補充）"
    return f"{title}\n{desc}"

obs1 = _extract_obs(_s1_lines, 0)
obs2 = _extract_obs(_s1_lines, 1)
obs3 = _extract_obs(_s1_lines, 2)

# ── Slide 2：排行榜重點觀察 ────────────────────────────────
_cnt_top3  = "、".join(f"{r['建案名稱']}({r['區']},{int(r['n'])}件)"
                        for _, r in rank_cnt.head(3).iterrows())
_unit_top3 = "、".join(f"{r['建案名稱']}({r['區']},{r['v']:.1f}萬/坪)"
                        for _, r in rank_unit.head(3).iterrows())
_tot_top3  = "、".join(f"{r['建案名稱']}({r['區']},{r['v']:.0f}萬)"
                        for _, r in rank_tot.head(3).iterrows())

_s2_prompt = f"""{_base}
成交量Top3建案：{_cnt_top3}
最高均單Top3建案：{_unit_top3}
最高總價Top3建案：{_tot_top3}
請針對以上三大排行榜寫2句重點觀察（共2句，每句不超過35字，不要編號或標點以外的符號）："""

s2_obs = ask_gemini(_s2_prompt)

# ── Slide 3：行政區分析（3組標題+說明）────────────────────
_dist_data = "\n".join(
    f"{r['行政區']}：{int(r['成交筆數'])}件，均單{r['均單']}萬/坪"
    for _, r in ds.head(10).iterrows()
)
_s3_prompt = f"""{_base}
行政區成交排行（前10）：
{_dist_data}
請輸出3組行政區市場觀察，格式嚴格如下（共6行，不要其他內容）：
標題1（6字以內）
說明1（30字以內，含具體數字）
標題2（6字以內）
說明2（30字以內，含具體數字）
標題3（6字以內）
說明3（30字以內，含具體數字）"""

_s3_raw   = ask_gemini(_s3_prompt)
_s3_lines = [l.strip() for l in _s3_raw.splitlines() if l.strip()]
def _s3_line(i): return _s3_lines[i] if i < len(_s3_lines) else "（待補充）"

s3_t1, s3_d1 = _s3_line(0), _s3_line(1)
s3_t2, s3_d2 = _s3_line(2), _s3_line(3)
s3_t3, s3_d3 = _s3_line(4), _s3_line(5)

# ── Slide 4：建物型態說明 ──────────────────────────────────
_type_str = "、".join(f"{k}({v}件)" for k, v in type_cnt.items())
_s4_prompt = f"""{_base}
建物型態分佈：{_type_str}
各行政區量價關係（成交量 vs 均單前5區）：
{chr(10).join(f"{r['行政區']}：{int(r['成交筆數'])}件，均單{r['均單']}萬/坪" for _, r in ds.head(5).iterrows())}
請寫1段40字以內的說明，描述本月建物型態特徵與量價分佈重點："""

s4_obs = ask_gemini(_s4_prompt)

# ── Slide 5：趨勢判讀 + 本月總結 ──────────────────────────
_trend_lines = []
for d in top5:
    vals = []
    for m in months_asc:
        sub = df[(df["鄉鎮市區"] == d) & (df["成交年月"] == m)]
        v = round(float(sub["單價"].mean()), 1) if len(sub) > 0 and not sub["單價"].isna().all() else None
        vals.append(f"{m}:{v}萬/坪" if v else f"{m}:-")
    _trend_lines.append(f"{d}：{', '.join(vals)}")

_s5_prompt = f"""{_base}
前5行政區均單月度趨勢：
{chr(10).join(_trend_lines)}
總成交量趨勢：{', '.join(f"{m}:{len(df[df['成交年月']==m])}件" for m in months_asc)}

請輸出2段（格式嚴格如下，共4行，不要其他內容）：
市場判讀標題（6字以內）
市場判讀說明（40字以內，描述均單趨勢）
本月總結標題（6字以內）
本月總結說明（40字以內，描述量能與展望）"""

_s5_raw   = ask_gemini(_s5_prompt)
_s5_lines = [l.strip() for l in _s5_raw.splitlines() if l.strip()]
def _s5_line(i): return _s5_lines[i] if i < len(_s5_lines) else "（待補充）"

s5_t7, s5_d7   = _s5_line(0), _s5_line(1)   # 文字方塊 7
s5_t11, s5_d11 = _s5_line(2), _s5_line(3)   # 文字方塊 11

# ══════════════════════════════════════════════════════════════
# 文字更新工具
# ══════════════════════════════════════════════════════════════

def find(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def fill_tf(shape, text):
    """填入文字方塊，\\n 對應多段落，保留既有 run 格式。"""
    if not shape or not shape.has_text_frame:
        return
    tf    = shape.text_frame
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i < len(tf.paragraphs):
            para = tf.paragraphs[i]
            if para.runs:
                para.runs[0].text = line
                for r in para.runs[1:]:
                    r.text = ""
    for i in range(len(lines), len(tf.paragraphs)):
        for r in tf.paragraphs[i].runs:
            r.text = ""

def set_single_run(shape, text, color=None):
    if not shape or not shape.has_text_frame:
        return
    para = shape.text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        return
    runs[0].text = text
    if color:
        runs[0].font.color.rgb = color
    for r in runs[1:]:
        r.text = ""

def update_kpi_val(shape, number_text):
    if not shape or not shape.has_text_frame:
        return
    para = shape.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = number_text

def update_kpi_chg(shape, text, color):
    if not shape or not shape.has_text_frame:
        return
    para = shape.text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        return
    runs[0].text = text
    runs[0].font.color.rgb = color
    for r in runs[1:]:
        r.text = ""

def get_kpi_groups(slide):
    result = {}
    for shape in slide.shapes:
        if not hasattr(shape, 'shapes'):
            continue
        lbl = val = chg = None
        for child in shape.shapes:
            n = child.name
            if n == 'KPI_lbl_0' and child.has_text_frame:
                lbl = child
            elif 'KPI_val' in n:
                val = child
            elif 'KPI_chg' in n:
                chg = child
        if lbl:
            result[lbl.text_frame.text.strip()] = (val, chg)
    return result

# ══════════════════════════════════════════════════════════════
# 圖表更新工具
# ══════════════════════════════════════════════════════════════

def find_chart(slide, shape_name):
    for s in slide.shapes:
        if s.name == shape_name and hasattr(s, 'chart'):
            return s.chart
    return None

def update_chart_month(chart, month):
    if not chart.has_title:
        return
    try:
        for para in chart.chart_title.text_frame.paragraphs:
            for run in para.runs:
                if re.search(r'\d{3}S\d|\d{4}-\d{2}', run.text):
                    run.text = re.sub(r'\d{3}S\d|\d{4}-\d{2}', month, run.text)
    except Exception as e:
        print(f"  警告：更新圖表標題失敗 {e}")

def update_bar(slide, shape_name, categories, values):
    chart = find_chart(slide, shape_name)
    if not chart:
        print(f"  警告：找不到圖表 {shape_name}"); return
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart.replace_data(cd)
    update_chart_month(chart, latest_m)

def update_doughnut(slide, shape_name, categories, values):
    chart = find_chart(slide, shape_name)
    if not chart:
        print(f"  警告：找不到圖表 {shape_name}"); return
    cd = ChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart.replace_data(cd)
    update_chart_month(chart, latest_m)

def update_bubble(slide, shape_name, data_rows):
    chart = find_chart(slide, shape_name)
    if not chart:
        print(f"  警告：找不到圖表 {shape_name}"); return
    bcd = BubbleChartData()
    for label, x, y, size in data_rows:
        s = bcd.add_series(label)
        s.add_data_point(float(x), float(y), float(size))
    chart.replace_data(bcd)
    update_chart_month(chart, latest_m)

def update_line(slide, shape_name, categories, series_list):
    chart = find_chart(slide, shape_name)
    if not chart:
        print(f"  警告：找不到圖表 {shape_name}"); return
    cd = ChartData()
    cd.categories = categories
    for name, values in series_list:
        cd.add_series(name, values)
    chart.replace_data(cd)

# ══════════════════════════════════════════════════════════════
# 填入模板
# ══════════════════════════════════════════════════════════════
print("填入模板...")
prs = Presentation(TEMPLATE)

# ── Slide 1：封面 KPI + 三大觀察 ──────────────────────────────
s1 = prs.slides[0]

set_single_run(find(s1, "TXT_subtitle"),
    f"統計月份：{latest_m}　｜　報表生成：{now.year} 年 {now.month:02d} 月 {now.day:02d} 日")

kpi = get_kpi_groups(s1)
if "總成交筆數" in kpi:
    update_kpi_val(kpi["總成交筆數"][0], f"{total_tx:,}")
    update_kpi_chg(kpi["總成交筆數"][1], chg_str(tx_chg, tx_pct, " 件"), chg_color(tx_chg))
if "全市平均均單" in kpi:
    update_kpi_val(kpi["全市平均均單"][0], f"{avg_unit}")
    update_kpi_chg(kpi["全市平均均單"][1], chg_str(u_chg, u_pct, " 萬"), chg_color(u_chg))
if "成交量\n最高行政區" in kpi:
    update_kpi_val(kpi["成交量\n最高行政區"][0], top_dist)
    update_kpi_chg(kpi["成交量\n最高行政區"][1], f"{top_tx} 件", C_GRAY)
if "活躍\n建案數" in kpi:
    update_kpi_val(kpi["活躍\n建案數"][0], str(num_proj))
if "涵蓋\n行政區數" in kpi:
    update_kpi_val(kpi["涵蓋\n行政區數"][0], str(num_dist))

fill_tf(find(s1, "文字方塊 13"), obs1)
fill_tf(find(s1, "文字方塊 14"), obs2)
fill_tf(find(s1, "文字方塊 15"), obs3)

# ── Slide 2：三大排行榜 ───────────────────────────────────────
s2 = prs.slides[1]

update_bar(s2, "Chart 10",
           [f"{r['建案名稱']}({r['區']})" for _, r in rank_cnt.iterrows()],
           rank_cnt["n"].tolist())
update_bar(s2, "Chart 11",
           [f"{r['建案名稱']}({r['區']})" for _, r in rank_unit.iterrows()],
           rank_unit["v"].round(1).tolist())
update_bar(s2, "Chart 12",
           [f"{r['建案名稱']}({r['區']})" for _, r in rank_tot.iterrows()],
           rank_tot["v"].round(0).tolist())
fill_tf(find(s2, "文字方塊 6"), s2_obs)

# ── Slide 3：行政區分析 ───────────────────────────────────────
s3 = prs.slides[2]

top15 = ds.head(15)
update_bar(s3, "Chart 8",
           top15["行政區"].tolist(),
           top15["成交筆數"].tolist())
fill_tf(find(s3, "文字方塊 19"), s3_t1)
fill_tf(find(s3, "文字方塊 20"), s3_d1)
fill_tf(find(s3, "文字方塊 21"), s3_t2)
fill_tf(find(s3, "文字方塊 22"), s3_d2)
fill_tf(find(s3, "文字方塊 23"), s3_t3)
fill_tf(find(s3, "文字方塊 24"), s3_d3)

# ── Slide 4：建物型態 & 量價 ──────────────────────────────────
s4 = prs.slides[3]

update_doughnut(s4, "Chart 8", type_cnt.index.tolist(), type_cnt.values.tolist())
bubble_data = [(row["行政區"], row["成交筆數"], row["均單"], row["成交筆數"])
               for _, row in ds.iterrows() if pd.notna(row["均單"])]
update_bubble(s4, "Chart 9", bubble_data)
fill_tf(find(s4, "文字方塊 5"), s4_obs)

# ── Slide 5：月份趨勢 ─────────────────────────────────────────
s5 = prs.slides[4]

set_single_run(find(s5, "TXT_section"), "月份趨勢")

series = []
for d in top5:
    vals = []
    for m in months_asc:
        sub = df[(df["鄉鎮市區"] == d) & (df["成交年月"] == m)]
        vals.append(round(float(sub["單價"].mean()), 1)
                    if len(sub) > 0 and not sub["單價"].isna().all() else None)
    series.append((d, vals))
update_line(s5, "Chart 4", months_asc, series)

fill_tf(find(s5, "文字方塊 7"),  f"{s5_t7}\n{s5_d7}")
fill_tf(find(s5, "文字方塊 11"), f"{s5_t11}\n{s5_d11}")

# ── 儲存 PPTX ─────────────────────────────────────────────────
out = OUTPUT_DIR / f"台中市房地產市場月報_{now.strftime('%Y%m')}.pptx"
prs.save(out)
print(f"PPTX 已生成：{out}")

# ── 轉換 PDF（Mac：LibreOffice）──────────────────────────────
import shutil
pdf_out = out.with_suffix(".pdf")
soffice = (
    "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if Path("/Applications/LibreOffice.app").exists()
    else shutil.which("soffice")
)
if soffice:
    try:
        subprocess.run([
            soffice, "--headless", "--convert-to", "pdf",
            "--outdir", str(out.parent), str(out)
        ], check=True, timeout=120)
        print(f"PDF 已生成：{pdf_out}")
    except Exception as e:
        print(f"PDF 轉換失敗（PPTX 仍已保留）：{e}")
else:
    print("未安裝 LibreOffice，略過 PDF 轉換（PPTX 仍已保留）")
